"""
Convert a /create-keynote markdown output into a PowerPoint (.pptx) file.

Usage:
    python keynote_to_pptx.py <input.md> [output.pptx] [--images-dir <path>]

If output path is omitted, writes to the same directory as input with .pptx extension.
If --images-dir is provided, looks for slide_1.png, slide_2.png, etc. to embed in slides.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# -- Theme colors ----------------------------------------------------------

BG_COLOR = RGBColor(0x0F, 0x17, 0x2A)       # Dark navy
TITLE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)     # White
BULLET_COLOR = RGBColor(0xCC, 0xCC, 0xCC)    # Light gray
ACCENT_COLOR = RGBColor(0x4E, 0xA8, 0xDE)    # Blue accent
NOTE_COLOR = RGBColor(0x99, 0x99, 0x99)      # Muted gray
SUBTITLE_COLOR = RGBColor(0x88, 0xBB, 0xDD)  # Light blue


def _set_slide_bg(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text_box(slide, left, top, width, height, text, font_size, color,
                  bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def parse_keynote_md(text: str) -> dict:
    """Parse the markdown output of /create-keynote into structured data."""
    result = {
        "title": "",
        "flow": [],
        "takeaway": "",
        "slides": [],
    }

    # Extract title from first H1 or H2
    title_match = re.search(r'^#{1,2}\s+(?:PRESENTATION\s*\d*:\s*)?(.+)', text, re.MULTILINE)
    if title_match:
        result["title"] = title_match.group(1).strip()

    # Extract flow
    flow_match = re.search(r'## FLOW\s*\n(.*?)(?=\n## )', text, re.DOTALL)
    if flow_match:
        for line in flow_match.group(1).strip().split("\n"):
            line = line.strip()
            if line.startswith(("-", "*")):
                line = re.sub(r'^[-*]\s*\d*\.?\s*', '', line)
            elif re.match(r'^\d+\.', line):
                line = re.sub(r'^\d+\.\s*', '', line)
            if line:
                result["flow"].append(line)

    # Extract takeaway
    takeaway_match = re.search(r'## DESIRED TAKEAWAY\s*\n(.+?)(?=\n## |\n---|\Z)', text, re.DOTALL)
    if takeaway_match:
        result["takeaway"] = takeaway_match.group(1).strip()

    # Extract slides
    slide_pattern = re.compile(
        r'### Slide \d+:\s*(.+?)\n'
        r'(.*?)(?=### Slide \d+:|---\s*\n\*\*Source|---\s*$|\Z)',
        re.DOTALL
    )

    for match in slide_pattern.finditer(text):
        slide_title = match.group(1).strip()
        body = match.group(2)

        bullets = []
        bullet_match = re.search(r'\*\*Bullets\*\*:\s*\n(.*?)(?=\n\*\*)', body, re.DOTALL)
        if bullet_match:
            for line in bullet_match.group(1).strip().split("\n"):
                line = re.sub(r'^[-*]\s*', '', line.strip())
                if line:
                    bullets.append(line)

        image = ""
        image_match = re.search(r'\*\*Image\*\*:\s*(.+?)(?=\n\*\*|\n\n|\Z)', body, re.DOTALL)
        if image_match:
            image = image_match.group(1).strip()

        notes = []
        notes_match = re.search(r'\*\*Speaker notes\*\*:\s*\n(.*?)(?=\n---|\n###|\Z)', body, re.DOTALL)
        if notes_match:
            for line in notes_match.group(1).strip().split("\n"):
                line = re.sub(r'^[-*]\s*', '', line.strip())
                if line:
                    notes.append(line)

        result["slides"].append({
            "title": slide_title,
            "bullets": bullets,
            "image": image,
            "notes": notes,
        })

    return result


def build_pptx(data: dict, output_path: Path, images_dir: Path | None = None) -> None:
    """Build a PowerPoint file from parsed keynote data."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # Blank layout

    slide_width = prs.slide_width
    slide_height = prs.slide_height
    margin = Inches(0.8)
    content_width = slide_width - (margin * 2)

    # -- Title slide --
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, BG_COLOR)

    title_text = data.get("title", "Presentation")
    _add_text_box(slide, margin, Inches(2.0), content_width, Inches(1.5),
                  title_text, 44, TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)

    if data.get("takeaway"):
        _add_text_box(slide, Inches(1.5), Inches(4.0), slide_width - Inches(3), Inches(1.0),
                      data["takeaway"], 18, SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    # -- Content slides --
    for s in data["slides"]:
        slide = prs.slides.add_slide(blank_layout)
        _set_slide_bg(slide, BG_COLOR)

        # Title
        _add_text_box(slide, margin, Inches(0.5), content_width, Inches(1.0),
                      s["title"], 36, TITLE_COLOR, bold=True)

        # Accent line under title
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            margin, Inches(1.3), Inches(2.0), Pt(3)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_COLOR
        shape.line.fill.background()

        # Bullets
        if s["bullets"]:
            txBox = slide.shapes.add_textbox(
                margin, Inches(1.8), content_width, Inches(4.0)
            )
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, bullet in enumerate(s["bullets"]):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(20)
                p.font.color.rgb = BULLET_COLOR
                p.space_after = Pt(12)
                p.level = 0

        # Image: embed file if available, otherwise show text description
        slide_num = data["slides"].index(s) + 1
        image_path = None
        if images_dir:
            for ext in (".png", ".jpg", ".jpeg", ".webp"):
                candidate = images_dir / f"slide_{slide_num}{ext}"
                if candidate.exists():
                    image_path = candidate
                    break

        if image_path:
            # Place image on the right side of the slide
            img_left = slide_width - margin - Inches(5.0)
            img_top = Inches(1.8)
            img_height = Inches(4.0)
            slide.shapes.add_picture(
                str(image_path), img_left, img_top, height=img_height
            )
            # Narrow bullet text to left half when image is present
        elif s["image"]:
            _add_text_box(slide, margin, Inches(6.2), content_width, Inches(0.8),
                          "Visual: " + s["image"][:120], 10, NOTE_COLOR)

        # Speaker notes
        if s["notes"]:
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = "\n".join("- " + n for n in s["notes"])

    # -- Closing slide --
    slide = prs.slides.add_slide(blank_layout)
    _set_slide_bg(slide, BG_COLOR)
    _add_text_box(slide, margin, Inches(2.5), content_width, Inches(1.5),
                  "Thank You", 44, TITLE_COLOR, bold=True, alignment=PP_ALIGN.CENTER)
    if data.get("takeaway"):
        _add_text_box(slide, Inches(1.5), Inches(4.0), slide_width - Inches(3), Inches(1.0),
                      data["takeaway"], 16, SUBTITLE_COLOR, alignment=PP_ALIGN.CENTER)

    prs.save(str(output_path))


def main() -> None:
    args = [a for a in sys.argv[1:] if not a.startswith("--")]
    images_dir = None
    for i, a in enumerate(sys.argv[1:], 1):
        if a == "--images-dir" and i + 1 < len(sys.argv):
            images_dir = Path(sys.argv[i + 1])

    if not args:
        print("Usage: keynote_to_pptx.py <input.md> [output.pptx] [--images-dir <path>]")
        sys.exit(1)

    input_path = Path(args[0])
    if not input_path.exists():
        print(f"File not found: {input_path}")
        sys.exit(1)

    output_path = Path(args[1]) if len(args) >= 2 else input_path.with_suffix(".pptx")

    text = input_path.read_text(encoding="utf-8")
    data = parse_keynote_md(text)

    if not data["slides"]:
        print("No slides found in input. Expected '### Slide N: Title' format.")
        sys.exit(1)

    build_pptx(data, output_path, images_dir)
    img_note = f" (with images from {images_dir})" if images_dir else ""
    print(f"Generated {len(data['slides'])} slides{img_note} -> {output_path}")


if __name__ == "__main__":
    main()
